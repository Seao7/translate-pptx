import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from googletrans import Translator
import copy
import io
import time
from datetime import datetime

def detect_language(text):
    """Detect the language of the input text."""
    translator = Translator()
    try:
        detection = translator.detect(text)
        return detection.lang
    except Exception as e:
        print(f"Language detection failed for: {text} - {e}")
        return 'unknown'

def translate_text(text, target_lang='ja', source_lang='auto'):
    """Translate text using googletrans with retry logic."""
    translator = Translator()
    max_retries = 3
    
    for attempt in range(max_retries):
        try:
            translated = translator.translate(text, src=source_lang, dest=target_lang)
            return translated.text
        except Exception as e:
            print(f"Translation attempt {attempt + 1} failed for: {text} - {e}")
            if attempt < max_retries - 1:
                time.sleep(1)  # Wait before retry
            else:
                return text  # fallback to original if all attempts fail

def copy_shape(source_shape, target_slide):
    """Copy a shape from source slide to target slide, preserving images and other content."""
    try:
        # Get the shape element
        shape_element = copy.deepcopy(source_shape._element)
        
        # Add the shape element to target slide
        target_slide._element.cSld.spTree.append(shape_element)
        
        # If it's an image, we need to copy the relationship
        if hasattr(source_shape, 'image'):
            # This is a picture shape
            image_part = source_shape.image.blob
            # Add image to target slide's relationships
            image_rId = target_slide.part.package.next_partname('/ppt/media/image%d.png')
            target_slide.part.package.get_or_add_image_part(image_part)
            
    except Exception as e:
        print(f"Error copying shape: {e}")

def copy_slide_with_relationships(source_slide, target_presentation):
    """Copy a slide with all its content including images, preserving relationships."""
    try:
        # Try to use the same layout as the source slide
        slide_layout = target_presentation.slide_layouts[0]  # Default to first layout
        
        # Try to find matching layout
        if hasattr(source_slide, 'slide_layout') and hasattr(source_slide.slide_layout, 'name'):
            layout_name = source_slide.slide_layout.name
            for layout in target_presentation.slide_layouts:
                if layout.name == layout_name:
                    slide_layout = layout
                    break
        
        # Create new slide
        new_slide = target_presentation.slides.add_slide(slide_layout)
        
        # Clear the new slide of default content
        shapes_to_remove = []
        for shape in new_slide.shapes:
            shapes_to_remove.append(shape)
        
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except:
                pass
        
        # Copy all shapes from source slide
        for shape in source_slide.shapes:
            copy_shape_comprehensive(shape, new_slide, source_slide, target_presentation)
        
        return new_slide
        
    except Exception as e:
        print(f"Error copying slide: {e}")
        # Fallback: create slide with basic layout
        slide_layout = target_presentation.slide_layouts[0]
        return target_presentation.slides.add_slide(slide_layout)

def copy_shape_comprehensive(source_shape, target_slide, source_slide, target_presentation):
    """Comprehensive shape copying including images, text, and formatting."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        if source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Handle image shapes
            copy_image_shape(source_shape, target_slide, target_presentation)
            
        elif source_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            # Handle text boxes
            copy_text_shape(source_shape, target_slide)
            
        elif source_shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            # Handle placeholders
            copy_placeholder_shape(source_shape, target_slide)
            
        elif source_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Handle grouped shapes
            copy_group_shape(source_shape, target_slide, source_slide, target_presentation)
            
        else:
            # Handle other shape types (rectangles, circles, etc.)
            copy_generic_shape(source_shape, target_slide)
            
    except Exception as e:
        print(f"Error copying shape type {getattr(source_shape, 'shape_type', 'unknown')}: {e}")

def copy_image_shape(source_shape, target_slide, target_presentation):
    """Copy image shape preserving the actual image."""
    try:
        # Get image data
        image_blob = source_shape.image.blob
        
        # Add image to target slide
        left = source_shape.left
        top = source_shape.top
        width = source_shape.width
        height = source_shape.height
        
        pic = target_slide.shapes.add_picture(
            io.BytesIO(image_blob), left, top, width, height
        )
        
        # Copy any additional formatting
        if hasattr(source_shape, 'rotation'):
            pic.rotation = source_shape.rotation
            
    except Exception as e:
        print(f"Error copying image: {e}")

def copy_text_shape(source_shape, target_slide):
    """Copy text box with formatting."""
    try:
        # Create text box
        left = source_shape.left
        top = source_shape.top
        width = source_shape.width
        height = source_shape.height
        
        textbox = target_slide.shapes.add_textbox(left, top, width, height)
        
        # Copy text content and formatting
        text_frame = textbox.text_frame
        text_frame.text = source_shape.text
        
        # Copy paragraph formatting
        for i, source_paragraph in enumerate(source_shape.text_frame.paragraphs):
            if i == 0:
                target_paragraph = text_frame.paragraphs[0]
            else:
                target_paragraph = text_frame.paragraphs.add()
            
            target_paragraph.text = source_paragraph.text
            
            # Copy paragraph-level formatting
            try:
                target_paragraph.alignment = source_paragraph.alignment
            except:
                pass
                
            # Copy run-level formatting
            for j, source_run in enumerate(source_paragraph.runs):
                if j == 0 and target_paragraph.runs:
                    target_run = target_paragraph.runs[0]
                elif target_paragraph.runs:
                    target_run = target_paragraph.runs.add()
                else:
                    continue
                    
                target_run.text = source_run.text
                
                try:
                    target_run.font.size = source_run.font.size
                    target_run.font.name = source_run.font.name
                    target_run.font.bold = source_run.font.bold
                    target_run.font.italic = source_run.font.italic
                    if source_run.font.color.rgb:
                        target_run.font.color.rgb = source_run.font.color.rgb
                except:
                    pass
                    
    except Exception as e:
        print(f"Error copying text shape: {e}")

def copy_placeholder_shape(source_shape, target_slide):
    """Copy placeholder content."""
    try:
        # Find corresponding placeholder in target slide
        if hasattr(source_shape, 'placeholder_format'):
            placeholder_type = source_shape.placeholder_format.type
            
            for target_shape in target_slide.shapes:
                if (hasattr(target_shape, 'placeholder_format') and 
                    target_shape.placeholder_format.type == placeholder_type):
                    
                    if hasattr(source_shape, 'text_frame') and source_shape.text_frame:
                        target_shape.text = source_shape.text
                    break
        else:
            # Fallback: treat as text box
            copy_text_shape(source_shape, target_slide)
            
    except Exception as e:
        print(f"Error copying placeholder: {e}")

def copy_group_shape(source_shape, target_slide, source_slide, target_presentation):
    """Copy grouped shapes."""
    try:
        # For grouped shapes, we need to copy each shape individually
        # since python-pptx doesn't have direct group copying
        for shape in source_shape.shapes:
            copy_shape_comprehensive(shape, target_slide, source_slide, target_presentation)
    except Exception as e:
        print(f"Error copying group: {e}")

def copy_generic_shape(source_shape, target_slide):
    """Copy other shape types (rectangles, etc.)."""
    try:
        # This is a simplified approach - for complex shapes you might need more specific handling
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        
        if hasattr(source_shape, 'auto_shape_type'):
            left = source_shape.left
            top = source_shape.top
            width = source_shape.width
            height = source_shape.height
            
            new_shape = target_slide.shapes.add_shape(
                source_shape.auto_shape_type, left, top, width, height
            )
            
            # Copy text if present
            if hasattr(source_shape, 'text_frame') and source_shape.text_frame:
                new_shape.text = source_shape.text
                
            # Copy basic formatting
            try:
                if source_shape.fill.solid():
                    new_shape.fill.solid()
                    new_shape.fill.fore_color.rgb = source_shape.fill.fore_color.rgb
            except:
                pass
                
    except Exception as e:
        print(f"Error copying generic shape: {e}")

def translate_pptx_standard(input_pptx_file, target_lang='ja', source_lang='auto'):
    """Standard translation - replaces original text with translated text."""
    prs = Presentation(input_pptx_file)
    new_prs = copy.deepcopy(prs)  # deep copy to preserve original
    
    # Count total number of text runs for progress bar
    total_runs = 0
    text_runs = []
    
    for slide in new_prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():  # Only count non-empty text
                            total_runs += 1
                            text_runs.append(run)
    
    if total_runs == 0:
        st.warning("No text content found in the presentation.")
        return None
    
    progress = st.progress(0)
    status_text = st.empty()
    
    translated_count = 0
    skipped_count = 0
    
    for i, run in enumerate(text_runs):
        original_text = run.text.strip()
        
        # Skip very short text or numbers/symbols only
        if len(original_text) <= 2 or original_text.isdigit():
            skipped_count += 1
        else:
            status_text.text(f"Translating: {original_text[:50]}...")
            translated_text = translate_text(original_text, target_lang, source_lang)
            run.text = translated_text
            translated_count += 1
            
            # Small delay to avoid hitting API rate limits
            time.sleep(0.1)
        
        progress.progress((i + 1) / total_runs)
    
    status_text.text(f"Translation complete! Translated: {translated_count}, Skipped: {skipped_count}")
    
    output = io.BytesIO()
    new_prs.save(output)
    output.seek(0)
    return output

def translate_pptx_bilingual(input_pptx_file, target_lang='ja', source_lang='auto'):
    """Bilingual translation - creates alternating original and translated slides preserving all content including images."""
    
    # Reset file pointer
    input_pptx_file.seek(0)
    original_prs = Presentation(input_pptx_file)
    
    # Create a new presentation using the original as template
    bilingual_prs = Presentation()
    
    # Remove default slide
    if len(bilingual_prs.slides) > 0:
        rId = bilingual_prs.slides._sldIdLst[0].rId
        bilingual_prs.part.drop_rel(rId)
        del bilingual_prs.slides._sldIdLst[0]
    
    # Copy slide layouts from original presentation if possible
    # (This is limited by python-pptx, so we'll work with default layouts)
    
    st.info("Creating bilingual presentation with image preservation...")
    
    total_slides = len(original_prs.slides)
    progress = st.progress(0)
    status_text = st.empty()
    
    for i in range(total_slides):
        status_text.text(f"Processing slide pair {i + 1} of {total_slides}...")
        
        original_slide = original_prs.slides[i]
        
        # Step 1: Add original slide
        status_text.text(f"Adding original slide {i + 1}...")
        original_copied = copy_slide_with_relationships(original_slide, bilingual_prs)
        
        # Step 2: Create and add translated slide
        status_text.text(f"Creating translated slide {i + 1}...")
        translated_copied = copy_slide_with_relationships(original_slide, bilingual_prs)
        
        # Step 3: Translate text in the translated slide
        translated_count = 0
        for shape in translated_copied.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text.strip()
                        if len(original_text) > 2 and not original_text.isdigit():
                            translated_text = translate_text(original_text, target_lang, source_lang)
                            run.text = translated_text
                            translated_count += 1
                            time.sleep(0.05)  # Small delay for API limits
        
        progress.progress((i + 1) / total_slides)
    
    status_text.text(f"Bilingual presentation created with {total_slides * 2} slides (including images)!")
    
    output = io.BytesIO()
    bilingual_prs.save(output)
    output.seek(0)
    return output

def get_language_name(lang_code):
    """Get language name from code."""
    languages = {
        'en': 'English',
        'ja': 'Japanese',
        'es': 'Spanish',
        'fr': 'French',
        'de': 'German',
        'zh': 'Chinese',
        'ko': 'Korean',
        'auto': 'Auto-detect'
    }
    return languages.get(lang_code, lang_code)

# ---- Streamlit App ----
st.set_page_config(
    page_title="PPTX Translator",
    page_icon="🌐",
    layout="wide"
)

st.title("🌐 PPTX Language Translator")
st.markdown("Translate PowerPoint presentations between different languages with ease!")

# Sidebar for settings
st.sidebar.header("Translation Settings")

# Translation mode selection
translation_mode = st.sidebar.radio(
    "Translation Mode:",
    options=["standard", "bilingual"],
    format_func=lambda x: "Standard (Replace original)" if x == "standard" else "Bilingual (Original + Translation)",
    index=0
)

st.sidebar.markdown("---")

# Language selection
source_lang = st.sidebar.selectbox(
    "Source Language:",
    options=['auto', 'en', 'ja', 'es', 'fr', 'de', 'zh', 'ko'],
    format_func=get_language_name,
    index=0
)

target_lang = st.sidebar.selectbox(
    "Target Language:",
    options=['ja', 'en', 'es', 'fr', 'de', 'zh', 'ko'],
    format_func=get_language_name,
    index=0 if source_lang != 'ja' else 1
)

# Quick translation buttons
st.sidebar.markdown("### Quick Options")
col1, col2 = st.sidebar.columns(2)
with col1:
    if st.button("🇺🇸→🇯🇵", help="English to Japanese"):
        source_lang = 'en'
        target_lang = 'ja'
        st.rerun()

with col2:
    if st.button("🇯🇵→🇺🇸", help="Japanese to English"):
        source_lang = 'ja'
        target_lang = 'en'
        st.rerun()

# Main content area
col1, col2 = st.columns([2, 1])

with col1:
    mode_description = "Standard Mode: Replaces original text with translations" if translation_mode == "standard" else "Bilingual Mode: Creates alternating original and translated slides"
    st.markdown(f"### {mode_description}")
    st.markdown(f"**Translation: {get_language_name(source_lang)} → {get_language_name(target_lang)}**")
    
    uploaded_file = st.file_uploader(
        "Upload a PPTX file",
        type=["pptx"],
        help="Select a PowerPoint presentation file to translate"
    )
    
    if uploaded_file is not None:
        st.success(f"✅ File uploaded: {uploaded_file.name}")
        
        # File info
        file_size = len(uploaded_file.getvalue()) / 1024 / 1024  # MB
        st.info(f"File size: {file_size:.2f} MB")
        
        if translation_mode == "bilingual":
            st.info("📄 Bilingual mode will create a presentation with double the slides (original → translated → original → translated...)")
            st.success("🖼️ Images and formatting will be preserved in bilingual mode!")
        
        # Validation
        if source_lang == target_lang:
            st.error("⚠️ Source and target languages cannot be the same!")
        else:
            if st.button("🚀 Start Translation", type="primary"):
                start_time = time.time()
                
                with st.spinner(f"{'Creating bilingual' if translation_mode == 'bilingual' else 'Translating'} presentation... This may take a few minutes."):
                    if translation_mode == "standard":
                        translated_pptx = translate_pptx_standard(uploaded_file, target_lang, source_lang)
                    else:
                        translated_pptx = translate_pptx_bilingual(uploaded_file, target_lang, source_lang)
                
                if translated_pptx:
                    end_time = time.time()
                    duration = end_time - start_time
                    
                    st.success(f"🎉 {'Bilingual presentation' if translation_mode == 'bilingual' else 'Translation'} completed in {duration:.1f} seconds!")
                    
                    # Generate filename with timestamp
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    mode_suffix = "bilingual" if translation_mode == "bilingual" else "translated"
                    filename = f"{mode_suffix}_{get_language_name(source_lang).lower()}_to_{get_language_name(target_lang).lower()}_{timestamp}.pptx"
                    
                    st.download_button(
                        label="⬇️ Download Translated PPTX",
                        data=translated_pptx,
                        file_name=filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        type="primary"
                    )
                else:
                    st.error("❌ Translation failed. Please try again.")

with col2:
    st.markdown("### ℹ️ How to Use")
    st.markdown("""
    1. **Choose translation mode**:
       - Standard: Replaces text
       - Bilingual: Alternating slides
    2. **Select languages** in the sidebar
    3. **Upload** your PPTX file
    4. **Click translate** and wait
    5. **Download** the result
    
    ### 🔧 Features
    - **Two translation modes**
    - **Bidirectional translation** (EN↔JP and more)
    - **Auto language detection**
    - **Progress tracking**
    - **Retry mechanism** for failed translations
    - **Smart text filtering** (skips numbers, short text)
    - **🖼️ Image preservation** in bilingual mode
    
    ### 📝 Notes
    - **Bilingual mode** doubles slide count
    - **Images are now preserved** in bilingual mode
    - Large files may take several minutes
    - Basic formatting is preserved
    - Original file remains unchanged
    - Rate limiting prevents API errors
    """)

# Footer
st.markdown("---")
st.markdown("Built with ❤️ using Streamlit and Google Translate API")

# Warning about API limits
with st.expander("⚠️ Important Notes"):
    st.markdown("""
    - **API Limits**: Google Translate has usage limits. Large presentations may hit rate limits.
    - **Accuracy**: Machine translation may not be perfect. Review important content manually.
    - **Formatting**: Complex layouts might need minor adjustments after translation.
    - **Privacy**: Text is sent to Google Translate service for processing.
    - **Bilingual Mode**: Creates alternating slides (original → translated → original → translated...)
    - **🖼️ Image Preservation**: Images and most formatting are now preserved in bilingual mode.
    - **Performance**: Bilingual mode with images may take longer to process.
    """)
